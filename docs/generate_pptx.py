import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors
    c_bg = RGBColor(7, 17, 31)          # #07111F
    c_card = RGBColor(15, 28, 46)       # #0F1C2E
    c_primary = RGBColor(92, 200, 255)  # #5CC8FF
    c_emergency = RGBColor(255, 77, 79) # #FF4D4F
    c_success = RGBColor(34, 197, 94)   # #22C55E
    c_gold = RGBColor(255, 209, 102)    # #FFD166
    c_white = RGBColor(248, 250, 252)   # #F8FAFC
    c_muted = RGBColor(148, 163, 184)   # #94A3B8
    c_border = RGBColor(40, 60, 90)

    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_bg

    def add_header(slide, title_text, sub_text):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = c_white
        p.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = sub_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = c_muted
        p2.font.name = "Segoe UI"

    # =========================================================================
    # SLIDE 1: HERO
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)

    # Top Pill
    shp = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(3.5), Inches(0.4))
    shp.fill.solid()
    shp.fill.fore_color.rgb = c_card
    shp.line.color.rgb = c_primary
    p = shp.text_frame.paragraphs[0]
    p.text = "⚡ HELPGURU PLATFORM • OS 1.0"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = c_primary

    # Title
    tx = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(7.5), Inches(1.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Intelligent Emergency Response Engine"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = c_white

    p2 = tf.add_paragraph()
    p2.text = '"When Every Second Counts, Every Decision Matters."'
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = c_gold

    p3 = tf.add_paragraph()
    p3.text = "A national-scale, cloud-native event-driven backend platform engineered for near-optimal, explainable, and conflict-free disaster management."
    p3.font.size = Pt(14)
    p3.font.color.rgb = c_muted

    # Map Card
    shp = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(0.8), Inches(4.0), Inches(5.8))
    shp.fill.solid()
    shp.fill.fore_color.rgb = c_card
    shp.line.color.rgb = c_border
    p = shp.text_frame.paragraphs[0]
    p.text = "🇧🇩 BANGLADESH EMERGENCY GRID\n\n🔴 BARISHAL (CRITICAL)\n🔴 SYLHET (FLOOD)\n🔴 CHITTAGONG (HIGH)\n🟡 DHAKA (COMMAND HUB)\n🟡 RAJSHAHI (HAZARD)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = c_primary
    p.alignment = PP_ALIGN.CENTER

    # 3 Badges
    badges = [("< 100ms", "Dispatch Latency", c_primary), ("100%", "Conflict-Free Lock", c_success), ("99.99%", "High Availability", c_gold)]
    for idx, (b_val, b_lbl, b_col) in enumerate(badges):
        shp = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 2.5), Inches(5.8), Inches(2.3), Inches(0.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = c_card
        shp.line.color.rgb = c_border
        p = shp.text_frame.paragraphs[0]
        p.text = f"{b_val}\n{b_lbl}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = b_col
        p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: THE CHALLENGE
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header(slide2, "The National Disaster Challenge", "Disasters are not bottlenecked by lack of resources—they fail due to wrong, delayed, and uncoordinated decisions.")

    # Left Box
    shp = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(25, 15, 25)
    shp.line.color.rgb = c_emergency
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚨 Traditional Systems (Legacy)\n"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = c_emergency

    lines_old = [
        "• Dispatch Logic: Manual & Fragmented",
        "• Resource Allocation: Static & Double-Booked",
        "• Adaptability: Delayed & Rigid",
        "• Decision Basis: Intuition-based"
    ]
    for line in lines_old:
        p = tf.add_paragraph()
        p.text = f"\n{line}"
        p.font.size = Pt(15)
        p.font.color.rgb = c_white

    # Right Box
    shp = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.8))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(15, 30, 25)
    shp.line.color.rgb = c_success
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚡ HelpGuru Platform (Next-Gen AI)\n"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = c_success

    lines_new = [
        "• Dispatch Logic: Automated Event Pipeline",
        "• Resource Allocation: Conflict-Free Locks (Redlock)",
        "• Adaptability: Near Real-Time Re-routing",
        "• Decision Basis: Explainable Multi-Factor AI"
    ]
    for line in lines_new:
        p = tf.add_paragraph()
        p.text = f"\n{line}"
        p.font.size = Pt(15)
        p.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 3: WORKFLOW PIPELINE
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_header(slide3, "Overall System Workflow", "End-to-end event stream from incident creation to monitoring feedback loop.")

    steps = [
        ("🚨 Incident Ingestion", "GeoJSON Intake"),
        ("⚡ Kafka Queue", "Event Pipeline"),
        ("🧠 Decision Engine", "Multi-Factor Solver"),
        ("🚑 Atomic Dispatch", "Redlock Lock"),
        ("📊 Telemetry Monitor", "Prometheus Feedback")
    ]

    for idx, (title, desc) in enumerate(steps):
        shp = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 2.4), Inches(2.8), Inches(2.2), Inches(2.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = c_card
        shp.line.color.rgb = c_gold if idx == 2 else c_primary
        tf = shp.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"STEP 0{idx+1}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = c_primary
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{title}\n"
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = c_white
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = c_muted
        p3.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 4: MICROSERVICES
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_header(slide4, "Microservice Architecture", "Decoupled, event-driven services communicating via internal secure network mesh.")

    svcs = [
        ("🌐 API Gateway", "Sole public ingress point & JWT OAuth2 authentication.", c_primary, Inches(0.8), Inches(2.0)),
        ("📡 Intake Service", "Validates GeoJSON & Streams to Kafka event bus.", c_primary, Inches(6.9), Inches(2.0)),
        ("🛰️ Resource Tracker", "GPS telemetry, vehicle status & hospital capacity.", RGBColor(168, 85, 247), Inches(0.8), Inches(3.8)),
        ("🧠 Decision Engine", "Multi-factor optimization constraint solver.", c_gold, Inches(6.9), Inches(3.8)),
    ]

    for title, desc, col, l, t in svcs:
        shp = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(5.6), Inches(1.5))
        shp.fill.solid()
        shp.fill.fore_color.rgb = c_card
        shp.line.color.rgb = col
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = col
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = c_white

    # Banner
    shp = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(15, 35, 55)
    shp.line.color.rgb = c_primary
    p = shp.text_frame.paragraphs[0]
    p.text = "🔐 Security Ingress: API Gateway serves as single ingress; all internal microservices communicate via mTLS internal mesh."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 5: HERO - DECISION ENGINE
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_header(slide5, "Explainable Multi-Factor Decision Engine", "How HelpGuru thinks: Synthesizing dynamic environmental variables into deterministic dispatch choices.")

    # Left Box
    shp = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8))
    shp.fill.solid()
    shp.fill.fore_color.rgb = c_card
    shp.line.color.rgb = c_gold
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🧮 Objective Score Function\n"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = c_gold
    vars_list = [
        "• 🚨 Incident Severity (Scale 1-10)",
        "• 👥 Affected Population Volume",
        "• ⏳ Time Sensitivity Window",
        "• 🗺️ Travel Time (PostGIS Route)",
        "• 🌧️ Weather & Road Hazard Index",
        "• 🏥 Real-Time Hospital Capacity"
    ]
    for v in vars_list:
        p = tf.add_paragraph()
        p.text = v
        p.font.size = Pt(15)
        p.font.color.rgb = c_white

    # Right Terminal Box
    shp = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.8))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(4, 9, 19)
    shp.line.color.rgb = c_primary
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SYSTEM LOG: DISPATCH OPTIMIZER\n"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_primary
    logs = [
        "[20:31:08] New Intake: Barishal ---> CRITICAL",
        "[20:31:10] Scanning PostGIS Grid ---> 14 Units",
        "[20:31:12] Risk Multiplier ---> Calculated",
        "[20:31:14] Redlock Acquisition ---> Locked",
        "[20:31:15] Dispatch Command ---> COMPLETE"
    ]
    for log in logs:
        p = tf.add_paragraph()
        p.text = log
        p.font.size = Pt(14)
        p.font.color.rgb = c_success

    # =========================================================================
    # SLIDE 6: DYNAMIC RE-OPTIMIZATION
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header(slide6, "Dynamic Re-Optimization & Resiliency", "Continuous adaptation to changing disaster environments and system instance failures.")

    timeline = [
        ("20:31:08", "[INTAKE]", "Incident created in Barishal Region • Ambulance #A-12 Dispatched.", c_primary),
        ("20:31:12", "[RE-ROUTE]", "Road Blocked Detected: Highway N8 flooded. Recalculating route in <50ms...", c_emergency),
        ("20:31:16", "[RE-ALLOCATE]", "Hospital Capacity Reached: District Hospital full. Rerouting to General Hospital.", c_gold),
        ("20:31:19", "[FAILOVER]", "Node Crash Recovery: Engine Replica #2 auto-assumes state via Kafka Event Replay.", c_success),
    ]

    for idx, (tm, tag, msg, col) in enumerate(timeline):
        shp = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0 + idx * 1.2), Inches(11.7), Inches(1.0))
        shp.fill.solid()
        shp.fill.fore_color.rgb = c_card
        shp.line.color.rgb = col
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{tm}  {tag}  {msg}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 7: DASHBOARD
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header(slide7, "Scalability, Data & Security Dashboard", "Live system operational metrics across all 13 core Architectural Categories.")

    dash_rows = [
        ("🟢 Horizontal Scalability", "Kubernetes Horizontal Pod Autoscaler (HPA) • Stateless Services"),
        ("🟢 Polyglot Data Layer", "PostgreSQL (Transactional ACID) + PostGIS (Geospatial R-Tree Indexing)"),
        ("🟢 Distributed Caching & Locks", "Redis Cluster • Sub-millisecond Lookups & Redlock Race-Condition Prevention"),
        ("🟢 Security & Governance", "Gateway-Only Ingress • Role-Based Access Control (RBAC) • E2E TLS 1.3")
    ]

    for idx, (title, val) in enumerate(dash_rows):
        shp = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0 + idx * 1.2), Inches(11.7), Inches(1.0))
        shp.fill.solid()
        shp.fill.fore_color.rgb = c_card
        shp.line.color.rgb = c_border
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = c_white
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(14)
        p2.font.color.rgb = c_muted

    # =========================================================================
    # SLIDE 8: IMPACT
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_header(slide8, "National Emergency Impact", "Transforming disaster management through high-throughput backend engineering.")

    impacts = [
        ("⏱️", "< 100ms", "Faster Response", "Sub-second automated dispatch reduces critical response windows.", c_primary),
        ("🔒", "100%", "Zero Conflicts", "Distributed locks guarantee zero resource double-booking.", c_success),
        ("🛡️", "99.99%", "High Availability", "Kafka event replay keeps operations online during regional outages.", c_gold),
        ("📈", "MAX", "Resource Use", "Dynamic re-allocation ensures no vehicle or facility remains idle.", RGBColor(168, 85, 247))
    ]

    for idx, (icon, val, lbl, desc, col) in enumerate(impacts):
        shp = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 3.0), Inches(2.0), Inches(2.7), Inches(4.8))
        shp.fill.solid()
        shp.fill.fore_color.rgb = c_card
        shp.line.color.rgb = c_border
        tf = shp.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{icon}\n\n{val}\n"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"{lbl}\n"
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = c_white
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(13)
        p3.font.color.rgb = c_muted
        p3.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 9: CLOSING
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9)

    tx = slide9.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.0))
    tf = tx.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "⚡ HelpGuru\n"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = c_primary
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = '"When Every Second Counts, Every Decision Matters."\n'
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = c_gold
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "A resilient, scalable, and explainable emergency management backbone designed to protect lives nationwide."
    p3.font.size = Pt(16)
    p3.font.color.rgb = c_muted
    p3.alignment = PP_ALIGN.CENTER

    shp = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(5.8), Inches(2.5), Inches(0.6))
    shp.fill.solid()
    shp.fill.fore_color.rgb = c_card
    shp.line.color.rgb = c_primary
    p = shp.text_frame.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = c_primary
    p.alignment = PP_ALIGN.CENTER

    # Save Presentation
    output_path = "/Users/mushfiq/Desktop/Hackathons !/HelpGuru_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
